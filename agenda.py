import os
from pptx import Presentation
import logging
import argparse

# add file logging
logging.basicConfig(filename='process.log', level=logging.INFO)

# parse the argument and extract the deck id
parser = argparse.ArgumentParser()
parser.add_argument('deck', nargs=1, help='deck id like 2024_07')
parser.add_argument('--out', type=str, default=f"agenda.txt", help="AOAI api key")
args = parser.parse_args()

deck_id = args.deck[0]

base_dir = os.getcwd()
filename = os.path.join(base_dir, f"Source\\{deck_id} - Azure-Technical Update Briefing.pptx")

#make sure Source and Annotated folders exist under the base one
if(not os.path.exists("Source")):
    os.mkdir("Source")
if(not os.path.exists("Annotated")):
    os.mkdir("Annotated")

if(not os.path.exists(filename)):
    print(f"File {filename} does not exist")
    exit(1)

prs = Presentation(filename)
with open(args.out, "w") as f:
    try:
        for slide in prs.slides:
            title = slide.shapes.title
            if(title):
                try:
                    print(title.text)
                    f.write(title.text)
                    f.write('\n')
                except:
                    print("Strange character")
                    continue
            else:
                continue
    except Exception as e:
        logging.error(f"Error processing deck {deck_id}: {e}")