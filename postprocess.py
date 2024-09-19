import os
from pptx import Presentation
import re
import requests
import logging
from openai import AzureOpenAI
from bs4 import BeautifulSoup
import tiktoken
import argparse
from tenacity import (
    retry,
    stop_after_attempt,
    wait_random_exponential,
)

# parse the argument and extract the deck id, the source of customer data as well as OpenAI API endpoint and credentials
parser = argparse.ArgumentParser()
parser.add_argument('--endpoint', type=str, help="AOAI api endpoint")
parser.add_argument('--key', type=str, help="AOAI api key")
parser.add_argument('--deployment', type=str, default="gpt-4o", help="AOAI deployment")
parser.add_argument('--apiversion', type=str, default="2023-03-15-preview", help="AOAI API version")
parser.add_argument('--customer', type=str, help="consumption file such as 'post.txt'")
parser.add_argument('deck', nargs=1, help='deck id like 2024_07')
args = parser.parse_args()

deck_id = args.deck[0]
customer = args.customer
customer_name = os.path.splitext(customer)[0]

# add file logging
logging.basicConfig(filename='postprocess.log', level=logging.INFO)

base_dir = os.getcwd()
filename = os.path.join(base_dir, f"Annotated\\{deck_id} - Azure-Technical Update Briefing.pptx")
filename_out = os.path.join(base_dir, "Annotated\\{deck_id} - Azure-Technical Update Briefing {customer_name}.pptx")

#make sure Source and Annotated folders exist under the base one
if(not os.path.exists("Source")):
    os.mkdir("Source")
if(not os.path.exists("Annotated")):
    os.mkdir("Annotated")

client = AzureOpenAI(
    api_key=args.key,  
    api_version=args.apiversion,
    azure_endpoint = args.endpoint
    )

def extract_urls(notes_text):
    url_pattern = r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+'
    urls = re.findall(url_pattern, notes_text)
    return urls

def download_url(url):
    try:
        response = requests.get(url)
        if response.status_code == 200:
            soup = BeautifulSoup(response.content, 'html.parser')
            text = soup.get_text()
            # replace multiple occurrences of newlines with a single newline
            text = re.sub(r'[\r\n]+', '\n', text)
            return text
        else:
            print(f"Failed to download {url}")
            return None
    except:
        return None

def ensure_safe_content_len(content, max_tokens):
    encoding = tiktoken.get_encoding("cl100k_base")
    
    processed_content = content
    tokens = encoding.encode(content)
    while len(tokens) > max_tokens:
        cutoff = - 5 + len(processed_content) * max_tokens / len(tokens)
        processed_content = content[:int(cutoff)]
        tokens = encoding.encode(processed_content)
    return processed_content

@retry(wait=wait_random_exponential(min=1, max=60), stop=stop_after_attempt(5))
def relevance(context, summary, topic):
    response = client.chat.completions.create(
      model=args.deployment,
      messages=[
            {"role": "system", "content": f"You are a helpful assistant helping the user to assess"
             f"if the information about an update for a service in Azure is relevant for the company's business. "
             f"Here's a report of the customer's recent consumption of azure services.\n---\n{context}. "
             f"User will provide a summary of a technical announcement, please respond with true or false "
             f"if you think that the update is relevant for the customer. Apply common reasoning such as \n "
             f"- If there is a usage of Kubernetes/AKS, consumption of virtual machines can be a sign of them being part of a kubernetes Nodepool. \n "
             f"- Consumption of SignalR indicate implementation of realtime UI data visualization. "
             f"\nOnly respond with \"true\" or \"false\" on the first line and provide a short (2-3 sentences) reasoning on the next line"},
            {"role": "user", "content": f"{topic}\n{summary}"},
        ]
    )
    try:
        first_choice = response.choices[0]
        message = first_choice.message
        response = message.content
        # extract the first line of text and convert to lowercase
        response = response.split("\n")
        # join the rest of the text starting with index 1
        reasoning = "".join(response[1:]).strip()
        return response[0].lower().startswith("true"), reasoning
    except:
        return (True, "---")

def append_notes(slide, is_active, reasoning):
    notes_slide = slide.notes_slide
    is_active_text = "Relevant" if is_active else "Not relevant"
    notes_slide.notes_text_frame.text += f"\n[{is_active_text}]\n" + reasoning

def hide_slide(slide):
    slide.element.set("show", "0")  

prs = Presentation(filename)
with open(f"{customer}", "r") as f:
    context = f.read()

for slide in prs.slides:
    title_shape = slide.shapes.title
    title = ""
    if(title_shape):
        try:
            title = title_shape.text
        except:
            continue
    else:
        continue

    if(slide.has_notes_slide):
        notes_slide = slide.notes_slide
        if "Autogenerated summary" not in notes_slide.notes_text_frame.text:
            continue

        notes_text = notes_slide.notes_text_frame.text
        urls = extract_urls(notes_text)
        summary = ""
        # take only the last 2 urls
        urls = urls[-2:]

        for url in urls:
            try:
                page_content = download_url(url)
                if page_content:
                    summary += "---\n" + page_content
            except:
                logging.error("Failed to download page")
                logging.error(url)
                continue
        
        summary = ensure_safe_content_len(summary, 5000)

        #summary = notes_text.split("[Autogenerated summary]")[1].strip()

        (is_active, reasoning) = relevance(context, summary, title)
        print(f"{title}: {'Relevant' if is_active else 'Not relevant'} ({reasoning})")
        append_notes(slide, is_active, reasoning)
        if not is_active:
            hide_slide(slide)

prs.save(filename_out)